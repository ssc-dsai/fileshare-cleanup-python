# Ingestion/extractors.py
"""
Text extraction dispatcher for supported file types.
Uses proper libraries for .docx / .pptx instead of raw read.
Adds vision flag only once per file.
"""

import logging
from pathlib import Path

import fitz  # PyMuPDF for PDF
import pytesseract
from PIL import Image

logger = logging.getLogger(__name__)


def extract_text_from_file(file_path: str | Path) -> str:
    """Extract readable text from file based on extension. Vision flag added only once."""
    path = Path(file_path)
    if not path.is_file():
        raise FileNotFoundError(f"File not found: {path}")

    ext = path.suffix.lower()

    try:
        if ext == ".pdf":
            text = _extract_pdf(path)
        elif ext in {".jpg", ".jpeg", ".png", ".gif", ".tif", ".tiff"}:
            text = _extract_image(path)  # flag added here — only once for images
        elif ext in {".doc", ".docx"}:
            text = _extract_docx(path)
        elif ext in {".ppt", ".pptx"}:
            text = _extract_pptx(path)
        elif ext == ".txt":
            text = path.read_text(encoding="utf-8", errors="replace").strip()
        else:
            raise ValueError(f"Unsupported extension: {ext}")

        # For non-image files: add [VISION_FLAG: No] **only if missing**
        if ext not in {".jpg", ".jpeg", ".png", ".gif", ".tif", ".tiff"}:
            if "[VISION_FLAG:" not in text:
                text = "[VISION_FLAG: No]\n\n" + text

        return text

    except Exception as e:
        logger.error(f"Extraction failed {path.name}: {type(e).__name__} - {str(e)}")
        return "[VISION_FLAG: No]\n[EXTRACTION FAILED]"


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using PyMuPDF."""
    text_parts = []
    doc = fitz.open(path)
    for page in doc:
        text = page.get_text("text").strip()
        if text:
            text_parts.append(text)
    doc.close()
    return "\n\n".join(text_parts)


def _extract_image(path: Path) -> str:
    """Perform OCR on image and add vision flag **once**."""
    try:
        img = Image.open(path)
        text = pytesseract.image_to_string(img, lang="eng+fra").strip()
        return f"[VISION_FLAG: Yes]\n\n{text}"
    except Exception as e:
        logger.warning(f"OCR failed on image {path.name}: {e}")
        return "[VISION_FLAG: Yes]\n[OCR extraction failed]"


def _extract_docx(path: Path) -> str:
    """Extract text from .docx using python-docx."""
    try:
        from docx import Document
        doc = Document(path)
        return "\n\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.warning(f"docx extraction failed {path.name}: {e}")
        return "[DOCX EXTRACTION FAILED]"


def _extract_pptx(path: Path) -> str:
    """Extract text from .pptx using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.warning(f"pptx extraction failed {path.name}: {e}")
        return "[PPTX EXTRACTION FAILED]"